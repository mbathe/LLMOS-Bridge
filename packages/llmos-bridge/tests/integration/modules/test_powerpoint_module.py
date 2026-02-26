"""Integration tests — PowerPointModule against real temp .pptx files."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from llmos_bridge.modules.powerpoint import PowerPointModule


@pytest.fixture
def module() -> PowerPointModule:
    return PowerPointModule()


@pytest.fixture
def pptx_path(tmp_path: Path) -> Path:
    """Create a simple presentation with two slides."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # Blank layout

    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide1.shapes.title.text = "Slide One"
    slide1.placeholders[1].text = "Subtitle text"

    slide2 = prs.slides.add_slide(blank_layout)
    txBox = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    txBox.text_frame.text = "Slide Two content"

    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


@pytest.mark.integration
class TestPresentationLifecycle:
    async def test_create_presentation(self, module: PowerPointModule, tmp_path: Path) -> None:
        out_path = tmp_path / "new.pptx"
        result = await module._action_create_presentation({"output_path": str(out_path)})
        assert out_path.exists()
        assert "path" in result

    async def test_open_presentation(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_open_presentation({"path": str(pptx_path)})
        assert result["slide_count"] == 2
        assert "layout_names" in result

    async def test_open_nonexistent_raises(self, module: PowerPointModule, tmp_path: Path) -> None:
        with pytest.raises(FileNotFoundError):
            await module._action_open_presentation({"path": str(tmp_path / "ghost.pptx")})

    async def test_save_presentation(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_save_presentation({"path": str(pptx_path)})
        assert "saved_to" in result

    async def test_save_to_new_path(self, module: PowerPointModule, pptx_path: Path, tmp_path: Path) -> None:
        out_path = tmp_path / "copy.pptx"
        result = await module._action_save_presentation(
            {"path": str(pptx_path), "output_path": str(out_path)}
        )
        assert out_path.exists()
        assert "saved_to" in result

    async def test_get_presentation_info(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_get_presentation_info({"path": str(pptx_path)})
        assert result["slide_count"] == 2
        assert "slide_layouts" in result
        assert result["slide_width_cm"] > 0


@pytest.mark.integration
class TestSlideManagement:
    async def test_list_slides(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_list_slides({"path": str(pptx_path)})
        assert result["slide_count"] == 2
        assert len(result["slides"]) == 2
        assert result["slides"][0]["title"] == "Slide One"

    async def test_read_slide(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_read_slide({"path": str(pptx_path), "slide_index": 0})
        assert result["index"] == 0
        assert result["title"] == "Slide One"

    async def test_add_slide(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_add_slide(
            {"path": str(pptx_path), "layout_index": 6, "title": "New Slide"}
        )
        assert result["slide_count"] == 3

    async def test_delete_slide(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_delete_slide(
            {"path": str(pptx_path), "slide_index": 1}
        )
        assert result["deleted_index"] == 1
        assert result["slide_count"] == 1

    async def test_delete_slide_out_of_range(self, module: PowerPointModule, pptx_path: Path) -> None:
        with pytest.raises(IndexError):
            await module._action_delete_slide({"path": str(pptx_path), "slide_index": 99})

    async def test_reorder_slide(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_reorder_slide(
            {"path": str(pptx_path), "from_index": 0, "to_index": 1}
        )
        assert result["from_index"] == 0
        assert result["to_index"] == 1

    async def test_set_slide_title(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_set_slide_title(
            {"path": str(pptx_path), "slide_index": 0, "title": "Updated Title"}
        )
        assert result["title"] == "Updated Title"


@pytest.mark.integration
class TestShapeOperations:
    async def test_add_text_box(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_add_text_box(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "text": "Hello from textbox",
                "left": 2.0,
                "top": 2.0,
                "width": 10.0,
                "height": 3.0,
            }
        )
        assert result["text"] == "Hello from textbox"

    async def test_add_text_box_with_formatting(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_add_text_box(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "text": "Formatted",
                "left": 2.0,
                "top": 5.0,
                "width": 8.0,
                "height": 2.0,
                "bold": True,
                "font_size": 24,
                "font_color": "FF0000",
            }
        )
        assert result["text"] == "Formatted"

    async def test_add_slide_notes(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_add_slide_notes(
            {"path": str(pptx_path), "slide_index": 0, "notes": "Speaker notes here"}
        )
        assert "notes" in result or "slide_index" in result

    async def test_add_shape_rectangle(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_add_shape(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "shape_type": "rectangle",
                "left": 1.0,
                "top": 1.0,
                "width": 5.0,
                "height": 3.0,
            }
        )
        assert "shape_index" in result

    async def test_add_table(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_add_table(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "rows": 3,
                "cols": 3,
                "left": 1.0,
                "top": 3.0,
                "width": 20.0,
                "height": 6.0,
                "data": [["H1", "H2", "H3"], ["r1c1", "r1c2", "r1c3"], ["r2c1", "r2c2", "r2c3"]],
            }
        )
        assert result["rows"] == 3
        assert result["cols"] == 3

    async def test_add_chart(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_add_chart(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "chart_type": "bar",
                "left": 1.0,
                "top": 3.0,
                "width": 15.0,
                "height": 10.0,
                "data": {
                    "categories": ["Cat A", "Cat B", "Cat C"],
                    "series": [{"name": "Series 1", "values": [1, 2, 3]}],
                },
            }
        )
        assert "shape_index" in result


@pytest.mark.integration
class TestSlideBackground:
    async def test_set_background_solid_color(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_set_slide_background(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "color": "FF5733",
            }
        )
        assert result["modified_slides"] >= 1

    async def test_add_transition(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_add_transition(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "transition": "fade",
            }
        )
        assert result["transition"] == "fade"
        assert result["modified_slides"] >= 1


# ---------------------------------------------------------------------------
# Extended tests — duplicate, layout, format_shape, format_table_cell
# ---------------------------------------------------------------------------


@pytest.mark.integration
class TestDuplicateSlide:
    async def test_duplicate_slide(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_duplicate_slide(
            {"path": str(pptx_path), "slide_index": 0}
        )
        assert "new_index" in result
        assert result["slide_count"] == 3

    async def test_duplicate_slide_with_insert_after(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        result = await module._action_duplicate_slide(
            {"path": str(pptx_path), "slide_index": 0, "insert_after": 0}
        )
        assert result["new_index"] == 1
        assert result["slide_count"] == 3

    async def test_duplicate_slide_out_of_range_raises(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        with pytest.raises(IndexError):
            await module._action_duplicate_slide(
                {"path": str(pptx_path), "slide_index": 99}
            )


@pytest.mark.integration
class TestSetSlideLayout:
    async def test_set_slide_layout(self, module: PowerPointModule, pptx_path: Path) -> None:
        result = await module._action_set_slide_layout(
            {"path": str(pptx_path), "slide_index": 0, "layout_index": 1}
        )
        assert result["slide_index"] == 0
        assert "layout_name" in result

    async def test_set_slide_layout_out_of_range_slide_raises(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        with pytest.raises(IndexError):
            await module._action_set_slide_layout(
                {"path": str(pptx_path), "slide_index": 99, "layout_index": 0}
            )

    async def test_set_slide_layout_out_of_range_layout_raises(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        with pytest.raises(IndexError):
            await module._action_set_slide_layout(
                {"path": str(pptx_path), "slide_index": 0, "layout_index": 999}
            )


@pytest.mark.integration
class TestFormatShape:
    async def test_format_shape_fill_color(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        # Slide 1 (index 0) has shapes from the Title Slide layout
        result = await module._action_format_shape(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "shape_index": 0,
                "fill_color": "FF0000",
            }
        )
        assert result["slide_index"] == 0
        assert result["shape_index"] == 0

    async def test_format_shape_rotation(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        result = await module._action_format_shape(
            {
                "path": str(pptx_path),
                "slide_index": 1,
                "shape_index": 0,
                "rotation": 45.0,
            }
        )
        assert result["shape_index"] == 0

    async def test_format_shape_out_of_range_slide_raises(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        with pytest.raises(IndexError):
            await module._action_format_shape(
                {
                    "path": str(pptx_path),
                    "slide_index": 99,
                    "shape_index": 0,
                    "fill_color": "0000FF",
                }
            )

    async def test_format_shape_out_of_range_shape_raises(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        with pytest.raises(IndexError):
            await module._action_format_shape(
                {
                    "path": str(pptx_path),
                    "slide_index": 0,
                    "shape_index": 999,
                    "fill_color": "0000FF",
                }
            )


@pytest.fixture
def pptx_with_table(tmp_path: Path) -> Path:
    """Create a presentation with a table on slide 0."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    # Add a table shape directly via python-pptx
    from pptx.util import Inches
    shape = slide.shapes.add_table(3, 3, Inches(1), Inches(1), Inches(6), Inches(3))
    tbl = shape.table
    tbl.cell(0, 0).text = "H1"
    tbl.cell(0, 1).text = "H2"
    tbl.cell(0, 2).text = "H3"
    path = tmp_path / "table.pptx"
    prs.save(str(path))
    return path


@pytest.mark.integration
class TestFormatTableCell:
    async def test_format_table_cell(
        self, module: PowerPointModule, pptx_with_table: Path
    ) -> None:
        # Shape index 0 is the table (only shape on slide 0)
        result = await module._action_format_table_cell(
            {
                "path": str(pptx_with_table),
                "slide_index": 0,
                "shape_index": 0,
                "row": 0,
                "col": 0,
                "text": "Modified",
                "bold": True,
                "bg_color": "FFFF00",
            }
        )
        assert result["row"] == 0
        assert result["col"] == 0

    async def test_format_table_cell_non_table_raises(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        with pytest.raises(ValueError, match="not a table"):
            await module._action_format_table_cell(
                {
                    "path": str(pptx_path),
                    "slide_index": 0,
                    "shape_index": 0,
                    "row": 0,
                    "col": 0,
                }
            )

    async def test_format_table_cell_out_of_range_slide_raises(
        self, module: PowerPointModule, pptx_with_table: Path
    ) -> None:
        with pytest.raises(IndexError):
            await module._action_format_table_cell(
                {
                    "path": str(pptx_with_table),
                    "slide_index": 99,
                    "shape_index": 0,
                    "row": 0,
                    "col": 0,
                }
            )

    async def test_format_table_cell_out_of_range_shape_raises(
        self, module: PowerPointModule, pptx_with_table: Path
    ) -> None:
        with pytest.raises(IndexError):
            await module._action_format_table_cell(
                {
                    "path": str(pptx_with_table),
                    "slide_index": 0,
                    "shape_index": 99,
                    "row": 0,
                    "col": 0,
                }
            )

    async def test_format_table_cell_with_alignment_and_bold(
        self, module: PowerPointModule, pptx_with_table: Path
    ) -> None:
        # First put text in cell so there are runs to format
        await module._action_format_table_cell(
            {
                "path": str(pptx_with_table),
                "slide_index": 0,
                "shape_index": 0,
                "row": 0,
                "col": 0,
                "text": "Header",
            }
        )
        # Now apply formatting
        result = await module._action_format_table_cell(
            {
                "path": str(pptx_with_table),
                "slide_index": 0,
                "shape_index": 0,
                "row": 0,
                "col": 0,
                "text": "Formatted",
                "bold": True,
                "italic": True,
                "font_size": 12,
                "font_color": "FF0000",
                "alignment": "center",
                "bg_color": "FFFF00",
            }
        )
        assert result["row"] == 0
        assert result["col"] == 0


# ---------------------------------------------------------------------------
# More coverage — slide background, add_shape, format_shape branches
# ---------------------------------------------------------------------------


@pytest.mark.integration
class TestSlideBackgroundExtended:
    async def test_set_background_all_slides_when_no_index(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        """Apply background to all slides (no slide_index given)."""
        result = await module._action_set_slide_background(
            {
                "path": str(pptx_path),
                "color": "336699",
                # No slide_index — applies to all 2 slides
            }
        )
        assert result["modified_slides"] == 2

    async def test_set_background_out_of_range_slide_raises(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        with pytest.raises(IndexError):
            await module._action_set_slide_background(
                {
                    "path": str(pptx_path),
                    "slide_index": 99,
                    "color": "FF0000",
                }
            )

    async def test_set_background_gradient(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        result = await module._action_set_slide_background(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "gradient": {
                    "angle": 90,
                    "stops": [
                        {"position": 0.0, "color": "FF0000"},
                        {"position": 1.0, "color": "0000FF"},
                    ],
                },
            }
        )
        assert result["modified_slides"] >= 1


@pytest.mark.integration
class TestAddSlideExtended:
    async def test_add_slide_with_title(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        result = await module._action_add_slide(
            {
                "path": str(pptx_path),
                "layout_index": 0,
                "title": "My Slide Title",
            }
        )
        assert result["slide_index"] >= 0

    async def test_add_slide_with_position(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        result = await module._action_add_slide(
            {
                "path": str(pptx_path),
                "layout_index": 6,  # blank layout
                "position": 0,  # insert at the beginning
            }
        )
        assert result["slide_index"] >= 0


@pytest.mark.integration
class TestAddSlideNotesExtended:
    async def test_add_slide_notes_append(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        """Test appending notes to an existing note."""
        await module._action_add_slide_notes(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "notes": "First note.",
                "append": False,
            }
        )
        result = await module._action_add_slide_notes(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "notes": "Second note.",
                "append": True,
            }
        )
        assert result["notes_length"] > 0


@pytest.mark.integration
class TestAddShapeExtended:
    async def test_add_shape_with_fill_color(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        result = await module._action_add_shape(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "shape_type": "rectangle",
                "left": 1.0,
                "top": 1.0,
                "width": 5.0,
                "height": 3.0,
                "fill_color": "FF5733",
            }
        )
        assert "shape_index" in result

    async def test_add_shape_with_line_color_and_width(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        result = await module._action_add_shape(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "shape_type": "ellipse",
                "left": 2.0,
                "top": 2.0,
                "width": 4.0,
                "height": 2.0,
                "line_color": "000000",
                "line_width": 2,
            }
        )
        assert "shape_index" in result

    async def test_add_shape_with_text_and_font(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        result = await module._action_add_shape(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "shape_type": "rectangle",
                "left": 1.0,
                "top": 4.0,
                "width": 6.0,
                "height": 2.0,
                "text": "Hello Shape",
                "font_size": 14,
                "font_color": "FFFFFF",
            }
        )
        assert "shape_index" in result


@pytest.mark.integration
class TestFormatShapeExtended:
    async def test_format_shape_with_line_color_and_shadow(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        # Add a shape first
        await module._action_add_shape(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "shape_type": "rectangle",
                "left": 1.0,
                "top": 1.0,
                "width": 4.0,
                "height": 3.0,
                "fill_color": "FF5733",
            }
        )
        # Format it with line color and shadow
        result = await module._action_format_shape(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "shape_index": 0,
                "line_color": "0000FF",
                "line_width": 3,
                "shadow": True,
            }
        )
        assert result["slide_index"] == 0

    async def test_format_shape_remove_shadow(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        await module._action_add_shape(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "shape_type": "rectangle",
                "left": 2.0,
                "top": 2.0,
                "width": 4.0,
                "height": 3.0,
            }
        )
        # First add shadow, then remove it
        await module._action_format_shape(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "shape_index": 0,
                "shadow": True,
            }
        )
        result = await module._action_format_shape(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "shape_index": 0,
                "shadow": False,
            }
        )
        assert result["slide_index"] == 0

    async def test_format_shape_with_fill_color_and_transparency(
        self, module: PowerPointModule, pptx_path: Path
    ) -> None:
        await module._action_add_shape(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "shape_type": "rectangle",
                "left": 3.0,
                "top": 1.0,
                "width": 4.0,
                "height": 3.0,
                "fill_color": "FF5733",
            }
        )
        result = await module._action_format_shape(
            {
                "path": str(pptx_path),
                "slide_index": 0,
                "shape_index": 0,
                "fill_color": "00FF00",
                "transparency": 0.5,
            }
        )
        assert result["slide_index"] == 0
