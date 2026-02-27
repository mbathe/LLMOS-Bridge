"""PowerPoint module — Implementation.

Full python-pptx-powered presentation automation.  All blocking I/O is
offloaded to a thread pool via ``asyncio.to_thread``.  python-pptx is
imported lazily so the module can be imported without it installed — the
dependency check in ``_check_dependencies`` will raise ``ModuleLoadError``
at instantiation time if missing.
"""

from __future__ import annotations

import asyncio
import copy
import os
import shutil
import subprocess
import tempfile
import threading
from pathlib import Path
from typing import TYPE_CHECKING, Any

from llmos_bridge.exceptions import ModuleLoadError
from llmos_bridge.modules.base import BaseModule, Platform
from llmos_bridge.modules.manifest import ActionSpec, ModuleManifest, ParamSpec
from llmos_bridge.security.decorators import audit_trail, requires_permission, sensitive_action
from llmos_bridge.security.models import Permission, RiskLevel
from llmos_bridge.protocol.params.powerpoint import (
    AddChartParams,
    AddImageParams,
    AddShapeParams,
    AddSlideNotesParams,
    AddSlideParams,
    AddTableParams,
    AddTextBoxParams,
    AddTransitionParams,
    ApplyThemeParams,
    CreatePresentationParams,
    DeleteSlideParams,
    DuplicateSlideParams,
    ExportSlideAsImageParams,
    ExportToPdfParams,
    FormatShapeParams,
    FormatTableCellParams,
    GetPresentationInfoParams,
    ListSlidesParams,
    OpenPresentationParams,
    ReorderSlideParams,
    ReadSlideParams,
    SavePresentationParams,
    SetSlideBackgroundParams,
    SetSlideLayoutParams,
    SetSlideTitleParams,
)

if TYPE_CHECKING:
    pass


# ---------------------------------------------------------------------------
# Shape type mapping
# ---------------------------------------------------------------------------

_SHAPE_TYPE_MAP: dict[str, int] = {
    # Values correspond to pptx.enum.shapes.MSO_AUTO_SHAPE_TYPE members
    "rectangle": 1,           # RECTANGLE
    "rounded_rectangle": 5,   # ROUNDED_RECTANGLE
    "ellipse": 9,             # OVAL
    "triangle": 7,            # ISOSCELES_TRIANGLE
    "right_arrow": 13,        # RIGHT_ARROW
    "left_arrow": 34,         # LEFT_ARROW
    "up_arrow": 35,           # UP_ARROW
    "down_arrow": 36,         # DOWN_ARROW
    "pentagon": 56,           # PENTAGON
    "hexagon": 10,            # HEXAGON
    "star4": 189,             # STAR_4_POINT
    "star5": 92,              # STAR_5_POINT
    "star8": 191,             # STAR_8_POINT
    "callout": 100,           # RECTANGULAR_CALLOUT
    "cloud": 179,             # CLOUD
    "lightning": 138,         # LIGHTNING_BOLT
    "heart": 74,              # HEART
    "checkmark": 20,          # FLOWCHART_PROCESS (fallback — no native checkmark)
    "line": 20,               # fallback
    "connector": 20,          # fallback
}

# Chart type string → XL_CHART_TYPE member name
_CHART_TYPE_MAP: dict[str, str] = {
    "bar": "BAR_CLUSTERED",
    "col": "COLUMN_CLUSTERED",
    "line": "LINE",
    "pie": "PIE",
    "doughnut": "DOUGHNUT",
    "scatter": "XY_SCATTER",
    "area": "AREA",
    "bubble": "BUBBLE",
    "radar": "RADAR",
}

# Transition name → OOXML element tag (p:transition child element name)
_TRANSITION_MAP: dict[str, str] = {
    "none": "",
    "fade": "fade",
    "push": "push",
    "wipe": "wipe",
    "split": "split",
    "reveal": "reveal",
    "random": "random",
}

# Alignment string → PP_ALIGN member
_ALIGN_MAP: dict[str, int] = {
    "left": 1,     # PP_ALIGN.LEFT
    "center": 2,   # PP_ALIGN.CENTER
    "right": 3,    # PP_ALIGN.RIGHT
    "justify": 4,  # PP_ALIGN.JUSTIFY
}


class PowerPointModule(BaseModule):
    """Full-featured PowerPoint automation backed by python-pptx."""

    MODULE_ID = "powerpoint"
    VERSION = "1.0.0"
    SUPPORTED_PLATFORMS = [Platform.ALL]

    def __init__(self) -> None:
        self._prs_cache: dict[str, Any] = {}
        self._path_locks: dict[str, threading.Lock] = {}
        self._meta_lock = threading.Lock()
        super().__init__()

    def _get_path_lock(self, path: str) -> threading.Lock:
        """Return (or create) a per-file threading.Lock for concurrent access control."""
        resolved = str(Path(path).resolve())
        with self._meta_lock:
            if resolved not in self._path_locks:
                self._path_locks[resolved] = threading.Lock()
            return self._path_locks[resolved]

    # ------------------------------------------------------------------
    # Dependency check
    # ------------------------------------------------------------------

    def _check_dependencies(self) -> None:
        try:
            import pptx  # noqa: F401
        except ImportError as exc:
            raise ModuleLoadError(
                module_id=self.MODULE_ID,
                reason="python-pptx is not installed. Run: pip install python-pptx",
            ) from exc

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _get_prs(self, path: str) -> Any:
        """Return a cached Presentation or load from disk."""
        from pptx import Presentation

        if path not in self._prs_cache:
            self._prs_cache[path] = Presentation(path)
        return self._prs_cache[path]

    def _save_prs(self, path: str, prs: Any, output_path: str | None = None) -> str:
        """Save presentation, update cache, return final path."""
        dest = output_path or path
        prs.save(dest)
        # Keep cache in sync
        self._prs_cache[dest] = prs
        if dest != path and path in self._prs_cache:
            del self._prs_cache[path]
        return dest

    @staticmethod
    def _hex_to_rgb(hex_color: str) -> Any:
        """Parse 'FF0000' or '#FF0000' → RGBColor(255, 0, 0)."""
        from pptx.dml.color import RGBColor

        hex_color = hex_color.lstrip("#")
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)

    @staticmethod
    def _cm(val: float) -> Any:
        """Convert cm float to EMU via pptx.util.Cm."""
        from pptx.util import Cm
        return Cm(val)

    @staticmethod
    def _emu_to_cm(emu: int) -> float:
        """Convert EMU to centimetres (rounded to 2 dp)."""
        return round(emu / 360000, 2)

    @staticmethod
    def _shape_info(shape: Any) -> dict[str, Any]:
        """Return a serialisable dict describing a shape."""
        info: dict[str, Any] = {
            "name": shape.name,
            "shape_type": str(shape.shape_type),
            "left_cm": PowerPointModule._emu_to_cm(shape.left or 0),
            "top_cm": PowerPointModule._emu_to_cm(shape.top or 0),
            "width_cm": PowerPointModule._emu_to_cm(shape.width or 0),
            "height_cm": PowerPointModule._emu_to_cm(shape.height or 0),
        }
        if shape.has_text_frame:
            paragraphs_text = [p.text for p in shape.text_frame.paragraphs]
            info["text"] = "\n".join(paragraphs_text)
        return info

    @staticmethod
    def _get_slide_title(slide: Any) -> str | None:
        """Return the title placeholder text, or None."""
        from pptx.enum.shapes import PP_PLACEHOLDER

        for ph in slide.placeholders:
            if ph.placeholder_format.type in (
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
            ):
                return ph.text
        # Fallback: check title attribute
        try:
            return slide.shapes.title.text if slide.shapes.title else None
        except Exception:
            return None

    @staticmethod
    def _apply_run_formatting(
        run: Any,
        bold: bool | None = None,
        italic: bool | None = None,
        underline: bool | None = None,
        font_name: str | None = None,
        font_size: int | None = None,
        font_color: str | None = None,
    ) -> None:
        """Apply formatting properties to a text run."""
        from pptx.util import Pt

        font = run.font
        if bold is not None:
            font.bold = bold
        if italic is not None:
            font.italic = italic
        if underline is not None:
            font.underline = underline
        if font_name:
            font.name = font_name
        if font_size is not None:
            font.size = Pt(font_size)
        if font_color:
            font.color.rgb = PowerPointModule._hex_to_rgb(font_color)

    @staticmethod
    def _apply_paragraph_alignment(paragraph: Any, alignment: str) -> None:
        from pptx.enum.text import PP_ALIGN

        align_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }
        if alignment in align_map:
            paragraph.alignment = align_map[alignment]

    # ------------------------------------------------------------------
    # Lifecycle actions
    # ------------------------------------------------------------------

    @requires_permission(Permission.FILESYSTEM_WRITE, reason="Modifies PowerPoint presentation")
    async def _action_create_presentation(self, params: dict[str, Any]) -> dict[str, Any]:
        p = CreatePresentationParams.model_validate(params)
        return await asyncio.to_thread(self._sync_create_presentation, p)

    def _sync_create_presentation(self, p: CreatePresentationParams) -> dict[str, Any]:
        from pptx import Presentation
        from pptx.util import Cm

        with self._get_path_lock(p.output_path):
            prs = Presentation()
            if p.slide_width is not None:
                prs.slide_width = Cm(p.slide_width)
            if p.slide_height is not None:
                prs.slide_height = Cm(p.slide_height)

            # Optionally copy theme from an existing pptx
            if p.theme_path:
                self._copy_theme_xml(p.theme_path, prs)

            Path(p.output_path).parent.mkdir(parents=True, exist_ok=True)
            prs.save(p.output_path)
            self._prs_cache[p.output_path] = prs
            return {
                "path": p.output_path,
                "slide_width_cm": self._emu_to_cm(prs.slide_width),
                "slide_height_cm": self._emu_to_cm(prs.slide_height),
            }

    async def _action_open_presentation(self, params: dict[str, Any]) -> dict[str, Any]:
        p = OpenPresentationParams.model_validate(params)
        return await asyncio.to_thread(self._sync_open_presentation, p)

    def _sync_open_presentation(self, p: OpenPresentationParams) -> dict[str, Any]:
        from pptx import Presentation

        with self._get_path_lock(p.path):
            if not Path(p.path).exists():
                raise FileNotFoundError(f"Presentation not found: {p.path}")
            prs = Presentation(p.path)
            self._prs_cache[p.path] = prs
            layout_names = [layout.name for layout in prs.slide_layouts]
            return {
                "path": p.path,
                "slide_count": len(prs.slides),
                "slide_width_cm": self._emu_to_cm(prs.slide_width),
                "slide_height_cm": self._emu_to_cm(prs.slide_height),
                "layout_names": layout_names,
            }

    @audit_trail("standard")
    @requires_permission(Permission.FILESYSTEM_WRITE, reason="Modifies PowerPoint presentation")
    async def _action_save_presentation(self, params: dict[str, Any]) -> dict[str, Any]:
        p = SavePresentationParams.model_validate(params)
        return await asyncio.to_thread(self._sync_save_presentation, p)

    def _sync_save_presentation(self, p: SavePresentationParams) -> dict[str, Any]:
        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            dest = self._save_prs(p.path, prs, p.output_path)
            return {"saved_to": dest}

    async def _action_get_presentation_info(self, params: dict[str, Any]) -> dict[str, Any]:
        p = GetPresentationInfoParams.model_validate(params)
        return await asyncio.to_thread(self._sync_get_presentation_info, p)

    def _sync_get_presentation_info(self, p: GetPresentationInfoParams) -> dict[str, Any]:
        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            layout_names = [layout.name for layout in prs.slide_layouts]
            return {
                "path": p.path,
                "slide_count": len(prs.slides),
                "slide_width_cm": self._emu_to_cm(prs.slide_width),
                "slide_height_cm": self._emu_to_cm(prs.slide_height),
                "slide_layouts": layout_names,
            }

    # ------------------------------------------------------------------
    # Slide management
    # ------------------------------------------------------------------

    @requires_permission(Permission.FILESYSTEM_WRITE, reason="Modifies PowerPoint presentation")
    async def _action_add_slide(self, params: dict[str, Any]) -> dict[str, Any]:
        p = AddSlideParams.model_validate(params)
        return await asyncio.to_thread(self._sync_add_slide, p)

    def _sync_add_slide(self, p: AddSlideParams) -> dict[str, Any]:
        from pptx.enum.shapes import PP_PLACEHOLDER

        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            layout = prs.slide_layouts[p.layout_index]
            slide = prs.slides.add_slide(layout)

            # Set title placeholder if requested
            if p.title:
                for ph in slide.placeholders:
                    if ph.placeholder_format.type in (
                        PP_PLACEHOLDER.TITLE,
                        PP_PLACEHOLDER.CENTER_TITLE,
                    ):
                        ph.text = p.title
                        break

            # If a position was requested, move the slide XML
            if p.position is not None:
                slide_list = prs.slides._sldIdLst
                # The new slide is at the end — move it to requested position
                slide_elem = slide_list[-1]
                slide_list.remove(slide_elem)
                target_pos = max(0, min(p.position, len(slide_list)))
                slide_list.insert(target_pos, slide_elem)

            new_index = len(prs.slides) - 1 if p.position is None else p.position
            self._save_prs(p.path, prs)
            return {"slide_index": new_index, "slide_count": len(prs.slides)}

    @sensitive_action(RiskLevel.MEDIUM)
    @requires_permission(Permission.FILESYSTEM_WRITE, reason="Modifies PowerPoint presentation")
    async def _action_delete_slide(self, params: dict[str, Any]) -> dict[str, Any]:
        p = DeleteSlideParams.model_validate(params)
        return await asyncio.to_thread(self._sync_delete_slide, p)

    def _sync_delete_slide(self, p: DeleteSlideParams) -> dict[str, Any]:
        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            slides = prs.slides
            if p.slide_index >= len(slides):
                raise IndexError(f"Slide index {p.slide_index} out of range ({len(slides)} slides).")

            # Get the slide object and its relationship ID
            slide = slides[p.slide_index]
            # Find rId in the slide list XML
            xml_slides = prs.slides._sldIdLst
            slide_elem = xml_slides[p.slide_index]
            r_id = slide_elem.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )

            # Remove the slide from the XML list
            xml_slides.remove(slide_elem)

            # Remove the relationship from the presentation part
            if r_id:
                prs.part.drop_rel(r_id)

            self._save_prs(p.path, prs)
            return {"deleted_index": p.slide_index, "slide_count": len(prs.slides)}

    @requires_permission(Permission.FILESYSTEM_WRITE, reason="Modifies PowerPoint presentation")
    async def _action_duplicate_slide(self, params: dict[str, Any]) -> dict[str, Any]:
        p = DuplicateSlideParams.model_validate(params)
        return await asyncio.to_thread(self._sync_duplicate_slide, p)

    def _sync_duplicate_slide(self, p: DuplicateSlideParams) -> dict[str, Any]:
        from lxml import etree
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.oxml.ns import qn

        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            slides = prs.slides
            if p.slide_index >= len(slides):
                raise IndexError(f"Slide index {p.slide_index} out of range.")

            source_slide = slides[p.slide_index]

            # Deep-copy the XML tree of the source slide
            source_xml = copy.deepcopy(source_slide._element)

            # Add a new slide using the same layout
            layout = source_slide.slide_layout
            new_slide = prs.slides.add_slide(layout)

            # Replace the new slide's XML element content with the copied source
            # (the element is the root of the slide XML document, so we copy in-place)
            el = new_slide._element
            for child in list(el):
                el.remove(child)
            for attr, val in source_xml.attrib.items():
                el.set(attr, val)
            for child in source_xml:
                el.append(copy.deepcopy(child))

            # Move to requested position
            if p.insert_after is not None:
                slide_list = prs.slides._sldIdLst
                new_elem = slide_list[-1]
                slide_list.remove(new_elem)
                target = min(p.insert_after + 1, len(slide_list))
                slide_list.insert(target, new_elem)

            self._save_prs(p.path, prs)
            new_index = (p.insert_after + 1) if p.insert_after is not None else len(prs.slides) - 1
            return {"new_index": new_index, "slide_count": len(prs.slides)}

    @requires_permission(Permission.FILESYSTEM_WRITE, reason="Modifies PowerPoint presentation")
    async def _action_reorder_slide(self, params: dict[str, Any]) -> dict[str, Any]:
        p = ReorderSlideParams.model_validate(params)
        return await asyncio.to_thread(self._sync_reorder_slide, p)

    def _sync_reorder_slide(self, p: ReorderSlideParams) -> dict[str, Any]:
        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            slides = prs.slides
            count = len(slides)
            if p.from_index >= count or p.to_index >= count:
                raise IndexError(f"Index out of range (slide count={count}).")

            slide_list = prs.slides._sldIdLst
            slide_elem = slide_list[p.from_index]
            slide_list.remove(slide_elem)
            slide_list.insert(p.to_index, slide_elem)

            self._save_prs(p.path, prs)
            return {"from_index": p.from_index, "to_index": p.to_index}

    async def _action_list_slides(self, params: dict[str, Any]) -> dict[str, Any]:
        p = ListSlidesParams.model_validate(params)
        return await asyncio.to_thread(self._sync_list_slides, p)

    def _sync_list_slides(self, p: ListSlidesParams) -> dict[str, Any]:
        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            result = []
            for idx, slide in enumerate(prs.slides):
                title = self._get_slide_title(slide)
                notes_preview = ""
                try:
                    if slide.has_notes_slide:
                        notes_tf = slide.notes_slide.notes_text_frame
                        notes_text = notes_tf.text
                        notes_preview = notes_text[:100] + ("..." if len(notes_text) > 100 else "")
                except Exception:
                    pass
                result.append(
                    {
                        "index": idx,
                        "title": title,
                        "shape_count": len(slide.shapes),
                        "notes_preview": notes_preview,
                    }
                )
            return {"slides": result, "slide_count": len(prs.slides)}

    async def _action_read_slide(self, params: dict[str, Any]) -> dict[str, Any]:
        p = ReadSlideParams.model_validate(params)
        return await asyncio.to_thread(self._sync_read_slide, p)

    def _sync_read_slide(self, p: ReadSlideParams) -> dict[str, Any]:
        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            if p.slide_index >= len(prs.slides):
                raise IndexError(f"Slide index {p.slide_index} out of range.")

            slide = prs.slides[p.slide_index]
            result: dict[str, Any] = {
                "index": p.slide_index,
                "title": self._get_slide_title(slide),
            }

            if p.include_shapes:
                result["shapes"] = [self._shape_info(shape) for shape in slide.shapes]

            if p.include_notes:
                notes = ""
                try:
                    if slide.has_notes_slide:
                        notes = slide.notes_slide.notes_text_frame.text
                except Exception:
                    pass
                result["notes"] = notes

            return result

    async def _action_set_slide_layout(self, params: dict[str, Any]) -> dict[str, Any]:
        p = SetSlideLayoutParams.model_validate(params)
        return await asyncio.to_thread(self._sync_set_slide_layout, p)

    def _sync_set_slide_layout(self, p: SetSlideLayoutParams) -> dict[str, Any]:
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT

        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            if p.slide_index >= len(prs.slides):
                raise IndexError(f"Slide index {p.slide_index} out of range.")
            if p.layout_index >= len(prs.slide_layouts):
                raise IndexError(f"Layout index {p.layout_index} out of range.")

            slide = prs.slides[p.slide_index]
            new_layout = prs.slide_layouts[p.layout_index]

            # python-pptx does not expose a setter for slide_layout; update the
            # underlying OPC relationship directly.
            slide_part = slide.part
            layout_part = new_layout.part
            for rId, rel in list(slide_part.rels.items()):
                if rel.reltype == RT.SLIDE_LAYOUT:
                    slide_part.rels.pop(rId)
                    break
            slide_part.relate_to(layout_part, RT.SLIDE_LAYOUT)

            self._save_prs(p.path, prs)
            return {"slide_index": p.slide_index, "layout_name": new_layout.name}

    # ------------------------------------------------------------------
    # Text content
    # ------------------------------------------------------------------

    async def _action_set_slide_title(self, params: dict[str, Any]) -> dict[str, Any]:
        p = SetSlideTitleParams.model_validate(params)
        return await asyncio.to_thread(self._sync_set_slide_title, p)

    def _sync_set_slide_title(self, p: SetSlideTitleParams) -> dict[str, Any]:
        from pptx.enum.shapes import PP_PLACEHOLDER
        from pptx.util import Pt

        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            if p.slide_index >= len(prs.slides):
                raise IndexError(f"Slide index {p.slide_index} out of range.")

            slide = prs.slides[p.slide_index]
            title_ph = None

            for ph in slide.placeholders:
                if ph.placeholder_format.type in (
                    PP_PLACEHOLDER.TITLE,
                    PP_PLACEHOLDER.CENTER_TITLE,
                ):
                    title_ph = ph
                    break

            if title_ph is None:
                # Try the generic title shape
                title_ph = slide.shapes.title

            if title_ph is None:
                raise ValueError(f"Slide {p.slide_index} has no title placeholder.")

            title_ph.text = p.title
            if p.bold or p.font_size or p.font_color:
                for para in title_ph.text_frame.paragraphs:
                    for run in para.runs:
                        if p.bold:
                            run.font.bold = True
                        if p.font_size:
                            run.font.size = Pt(p.font_size)
                        if p.font_color:
                            run.font.color.rgb = self._hex_to_rgb(p.font_color)

            self._save_prs(p.path, prs)
            return {"slide_index": p.slide_index, "title": p.title}

    @requires_permission(Permission.FILESYSTEM_WRITE, reason="Modifies PowerPoint presentation")
    async def _action_add_text_box(self, params: dict[str, Any]) -> dict[str, Any]:
        p = AddTextBoxParams.model_validate(params)
        return await asyncio.to_thread(self._sync_add_text_box, p)

    def _sync_add_text_box(self, p: AddTextBoxParams) -> dict[str, Any]:
        from pptx.enum.text import MSO_ANCHOR
        from pptx.util import Cm, Pt

        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            if p.slide_index >= len(prs.slides):
                raise IndexError(f"Slide index {p.slide_index} out of range.")

            slide = prs.slides[p.slide_index]
            txBox = slide.shapes.add_textbox(Cm(p.left), Cm(p.top), Cm(p.width), Cm(p.height))
            tf = txBox.text_frame
            tf.word_wrap = p.word_wrap

            # Vertical alignment
            v_align_map = {
                "top": MSO_ANCHOR.TOP,
                "middle": MSO_ANCHOR.MIDDLE,
                "bottom": MSO_ANCHOR.BOTTOM,
            }
            tf.vertical_anchor = v_align_map.get(p.vertical_alignment, MSO_ANCHOR.TOP)

            # Background fill
            if p.bg_color:
                txBox.fill.solid()
                txBox.fill.fore_color.rgb = self._hex_to_rgb(p.bg_color)

            # Set text and formatting
            para = tf.paragraphs[0]
            self._apply_paragraph_alignment(para, p.alignment)
            run = para.add_run()
            run.text = p.text
            self._apply_run_formatting(
                run,
                bold=p.bold,
                italic=p.italic,
                underline=p.underline,
                font_name=p.font_name,
                font_size=p.font_size,
                font_color=p.font_color,
            )

            self._save_prs(p.path, prs)
            return {
                "slide_index": p.slide_index,
                "shape_index": len(slide.shapes) - 1,
                "text": p.text,
            }

    @requires_permission(Permission.FILESYSTEM_WRITE, reason="Modifies PowerPoint presentation")
    async def _action_add_slide_notes(self, params: dict[str, Any]) -> dict[str, Any]:
        p = AddSlideNotesParams.model_validate(params)
        return await asyncio.to_thread(self._sync_add_slide_notes, p)

    def _sync_add_slide_notes(self, p: AddSlideNotesParams) -> dict[str, Any]:
        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            if p.slide_index >= len(prs.slides):
                raise IndexError(f"Slide index {p.slide_index} out of range.")

            slide = prs.slides[p.slide_index]
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame

            if p.append and notes_tf.text:
                # Append to existing notes
                para = notes_tf.add_paragraph()
                para.text = p.notes
            else:
                # Replace
                notes_tf.text = p.notes

            self._save_prs(p.path, prs)
            return {"slide_index": p.slide_index, "notes_length": len(notes_tf.text)}

    # ------------------------------------------------------------------
    # Shapes
    # ------------------------------------------------------------------

    @requires_permission(Permission.FILESYSTEM_WRITE, reason="Modifies PowerPoint presentation")
    async def _action_add_shape(self, params: dict[str, Any]) -> dict[str, Any]:
        p = AddShapeParams.model_validate(params)
        return await asyncio.to_thread(self._sync_add_shape, p)

    def _sync_add_shape(self, p: AddShapeParams) -> dict[str, Any]:
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        from pptx.util import Cm, Pt

        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            if p.slide_index >= len(prs.slides):
                raise IndexError(f"Slide index {p.slide_index} out of range.")

            slide = prs.slides[p.slide_index]

            # Resolve shape type enum
            shape_type_val = _SHAPE_TYPE_MAP.get(p.shape_type, 1)
            # Look up in MSO_AUTO_SHAPE_TYPE by value
            auto_shape_type = None
            for member in MSO_AUTO_SHAPE_TYPE:
                if member.value == shape_type_val:
                    auto_shape_type = member
                    break
            if auto_shape_type is None:
                auto_shape_type = MSO_AUTO_SHAPE_TYPE.RECTANGLE

            shape = slide.shapes.add_shape(
                auto_shape_type,
                Cm(p.left),
                Cm(p.top),
                Cm(p.width),
                Cm(p.height),
            )

            # Fill
            if p.fill_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = self._hex_to_rgb(p.fill_color)
                if p.transparency is not None:
                    shape.fill.fore_color.theme_color = None  # ensure RGB mode
                    # Transparency in OOXML is 0–100000 (int), but python-pptx
                    # doesn't expose it directly — set via XML
                    from lxml import etree

                    solidFill = shape.fill._xPr.find(
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill"
                    )
                    if solidFill is not None:
                        srgb = solidFill.find(
                            "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"
                        )
                        if srgb is not None:
                            alpha = etree.SubElement(
                                srgb,
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha",
                            )
                            alpha.set("val", str(int((1 - p.transparency) * 100000)))
            else:
                shape.fill.background()

            # Line
            if p.line_color:
                shape.line.color.rgb = self._hex_to_rgb(p.line_color)
            if p.line_width is not None:
                from pptx.util import Pt
                shape.line.width = Pt(p.line_width)

            # Text
            if p.text:
                tf = shape.text_frame
                tf.word_wrap = True
                para = tf.paragraphs[0]
                run = para.add_run()
                run.text = p.text
                if p.font_size:
                    run.font.size = Pt(p.font_size)
                if p.font_color:
                    run.font.color.rgb = self._hex_to_rgb(p.font_color)

            self._save_prs(p.path, prs)
            return {
                "slide_index": p.slide_index,
                "shape_index": len(slide.shapes) - 1,
                "shape_type": p.shape_type,
            }

    @requires_permission(Permission.FILESYSTEM_WRITE, reason="Modifies PowerPoint presentation")
    async def _action_format_shape(self, params: dict[str, Any]) -> dict[str, Any]:
        p = FormatShapeParams.model_validate(params)
        return await asyncio.to_thread(self._sync_format_shape, p)

    def _sync_format_shape(self, p: FormatShapeParams) -> dict[str, Any]:
        from pptx.util import Pt

        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            if p.slide_index >= len(prs.slides):
                raise IndexError(f"Slide index {p.slide_index} out of range.")

            slide = prs.slides[p.slide_index]
            if p.shape_index >= len(slide.shapes):
                raise IndexError(f"Shape index {p.shape_index} out of range.")

            shape = slide.shapes[p.shape_index]

            if p.fill_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = self._hex_to_rgb(p.fill_color)

            if p.transparency is not None and shape.fill.type is not None:
                from lxml import etree

                try:
                    solidFill = shape.fill._xPr.find(
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill"
                    )
                    if solidFill is not None:
                        srgb = solidFill.find(
                            "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"
                        )
                        if srgb is not None:
                            # Remove existing alpha if present
                            for old in srgb.findall(
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha"
                            ):
                                srgb.remove(old)
                            alpha = etree.SubElement(
                                srgb,
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha",
                            )
                            alpha.set("val", str(int((1 - p.transparency) * 100000)))
                except Exception:
                    pass

            if p.line_color:
                shape.line.color.rgb = self._hex_to_rgb(p.line_color)

            if p.line_width is not None:
                shape.line.width = Pt(p.line_width)

            if p.rotation is not None:
                shape.rotation = p.rotation

            if p.shadow is not None:
                # python-pptx doesn't expose shadow natively; inject XML
                from lxml import etree

                spPr = shape._element.spPr
                ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
                effLst = spPr.find(f"{{{ns}}}effectLst")
                if p.shadow:
                    if effLst is None:
                        effLst = etree.SubElement(spPr, f"{{{ns}}}effectLst")
                    outerShdw = effLst.find(f"{{{ns}}}outerShdw")
                    if outerShdw is None:
                        outerShdw = etree.SubElement(effLst, f"{{{ns}}}outerShdw")
                    outerShdw.set("blurRad", "40000")
                    outerShdw.set("dist", "23000")
                    outerShdw.set("dir", "5400000")
                    outerShdw.set("rotWithShape", "0")
                    srgb = etree.SubElement(outerShdw, f"{{{ns}}}srgbClr")
                    srgb.set("val", "000000")
                    alpha = etree.SubElement(srgb, f"{{{ns}}}alpha")
                    alpha.set("val", "63000")
                else:
                    if effLst is not None:
                        spPr.remove(effLst)

            self._save_prs(p.path, prs)
            return {"slide_index": p.slide_index, "shape_index": p.shape_index}

    @requires_permission(Permission.FILESYSTEM_WRITE, reason="Modifies PowerPoint presentation")
    async def _action_add_image(self, params: dict[str, Any]) -> dict[str, Any]:
        p = AddImageParams.model_validate(params)
        return await asyncio.to_thread(self._sync_add_image, p)

    def _sync_add_image(self, p: AddImageParams) -> dict[str, Any]:
        from pptx.util import Cm

        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            if p.slide_index >= len(prs.slides):
                raise IndexError(f"Slide index {p.slide_index} out of range.")
            if not Path(p.image_path).exists():
                raise FileNotFoundError(f"Image not found: {p.image_path}")

            slide = prs.slides[p.slide_index]
            width = Cm(p.width) if p.width else None
            height = Cm(p.height) if p.height else None

            pic = slide.shapes.add_picture(p.image_path, Cm(p.left), Cm(p.top), width=width, height=height)

            self._save_prs(p.path, prs)
            return {
                "slide_index": p.slide_index,
                "shape_index": len(slide.shapes) - 1,
                "width_cm": self._emu_to_cm(pic.width),
                "height_cm": self._emu_to_cm(pic.height),
            }

    # ------------------------------------------------------------------
    # Charts
    # ------------------------------------------------------------------

    @requires_permission(Permission.FILESYSTEM_WRITE, reason="Modifies PowerPoint presentation")
    async def _action_add_chart(self, params: dict[str, Any]) -> dict[str, Any]:
        p = AddChartParams.model_validate(params)
        return await asyncio.to_thread(self._sync_add_chart, p)

    def _sync_add_chart(self, p: AddChartParams) -> dict[str, Any]:
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Cm

        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            if p.slide_index >= len(prs.slides):
                raise IndexError(f"Slide index {p.slide_index} out of range.")

            slide = prs.slides[p.slide_index]

            # Resolve chart type
            chart_type_name = _CHART_TYPE_MAP.get(p.chart_type, "COLUMN_CLUSTERED")
            chart_type_enum = getattr(XL_CHART_TYPE, chart_type_name)

            # Build ChartData
            chart_data = ChartData()
            categories = p.data.get("categories", [])
            chart_data.categories = categories

            for series in p.data.get("series", []):
                chart_data.add_series(series["name"], series["values"])

            graphic_frame = slide.shapes.add_chart(
                chart_type_enum,
                Cm(p.left),
                Cm(p.top),
                Cm(p.width),
                Cm(p.height),
                chart_data,
            )

            chart = graphic_frame.chart
            chart.chart_style = p.style
            chart.has_legend = p.has_legend

            if p.title:
                chart.has_title = True
                chart.chart_title.has_text_frame = True
                chart.chart_title.text_frame.text = p.title

            if p.has_data_labels:
                plot = chart.plots[0]
                plot.has_data_labels = True

            self._save_prs(p.path, prs)
            return {
                "slide_index": p.slide_index,
                "chart_type": p.chart_type,
                "shape_index": len(slide.shapes) - 1,
            }

    # ------------------------------------------------------------------
    # Tables
    # ------------------------------------------------------------------

    @requires_permission(Permission.FILESYSTEM_WRITE, reason="Modifies PowerPoint presentation")
    async def _action_add_table(self, params: dict[str, Any]) -> dict[str, Any]:
        p = AddTableParams.model_validate(params)
        return await asyncio.to_thread(self._sync_add_table, p)

    def _sync_add_table(self, p: AddTableParams) -> dict[str, Any]:
        from pptx.util import Cm, Pt

        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            if p.slide_index >= len(prs.slides):
                raise IndexError(f"Slide index {p.slide_index} out of range.")

            slide = prs.slides[p.slide_index]
            graphic_frame = slide.shapes.add_table(
                p.rows, p.cols, Cm(p.left), Cm(p.top), Cm(p.width), Cm(p.height)
            )
            table = graphic_frame.table

            header_bg = self._hex_to_rgb(p.header_bg_color)
            header_fg = self._hex_to_rgb(p.header_font_color)
            alt_bg = self._hex_to_rgb(p.alt_row_color) if p.alt_row_color else None

            for row_idx in range(p.rows):
                for col_idx in range(p.cols):
                    cell = table.cell(row_idx, col_idx)
                    # Set cell text from data if provided
                    if p.data and row_idx < len(p.data) and col_idx < len(p.data[row_idx]):
                        cell.text = str(p.data[row_idx][col_idx])

                    # Font size for all cells
                    tf = cell.text_frame
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(p.font_size)

                    # Header row styling
                    if p.has_header and row_idx == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = header_bg
                        for para in tf.paragraphs:
                            for run in para.runs:
                                run.font.color.rgb = header_fg
                                run.font.bold = True
                    elif alt_bg and row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = alt_bg

            self._save_prs(p.path, prs)
            return {
                "slide_index": p.slide_index,
                "shape_index": len(slide.shapes) - 1,
                "rows": p.rows,
                "cols": p.cols,
            }

    async def _action_format_table_cell(self, params: dict[str, Any]) -> dict[str, Any]:
        p = FormatTableCellParams.model_validate(params)
        return await asyncio.to_thread(self._sync_format_table_cell, p)

    def _sync_format_table_cell(self, p: FormatTableCellParams) -> dict[str, Any]:
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Pt

        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            if p.slide_index >= len(prs.slides):
                raise IndexError(f"Slide index {p.slide_index} out of range.")

            slide = prs.slides[p.slide_index]
            if p.shape_index >= len(slide.shapes):
                raise IndexError(f"Shape index {p.shape_index} out of range.")

            shape = slide.shapes[p.shape_index]
            if not shape.has_table:
                raise ValueError(f"Shape {p.shape_index} is not a table.")

            table = shape.table
            cell = table.cell(p.row, p.col)

            if p.text is not None:
                cell.text = p.text

            if p.bg_color:
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._hex_to_rgb(p.bg_color)

            # Apply formatting to runs
            align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
            for para in cell.text_frame.paragraphs:
                if p.alignment and p.alignment in align_map:
                    para.alignment = align_map[p.alignment]
                for run in para.runs:
                    if p.bold is not None:
                        run.font.bold = p.bold
                    if p.italic is not None:
                        run.font.italic = p.italic
                    if p.font_size is not None:
                        run.font.size = Pt(p.font_size)
                    if p.font_color:
                        run.font.color.rgb = self._hex_to_rgb(p.font_color)

            self._save_prs(p.path, prs)
            return {"slide_index": p.slide_index, "shape_index": p.shape_index, "row": p.row, "col": p.col}

    # ------------------------------------------------------------------
    # Background & theme
    # ------------------------------------------------------------------

    @requires_permission(Permission.FILESYSTEM_WRITE, reason="Modifies PowerPoint presentation")
    async def _action_set_slide_background(self, params: dict[str, Any]) -> dict[str, Any]:
        p = SetSlideBackgroundParams.model_validate(params)
        return await asyncio.to_thread(self._sync_set_slide_background, p)

    def _sync_set_slide_background(self, p: SetSlideBackgroundParams) -> dict[str, Any]:
        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)

            # Determine target slides
            if p.slide_index is not None:
                if p.slide_index >= len(prs.slides):
                    raise IndexError(f"Slide index {p.slide_index} out of range.")
                target_slides = [prs.slides[p.slide_index]]
            else:
                target_slides = list(prs.slides)

            modified = 0
            for slide in target_slides:
                background = slide.background
                fill = background.fill

                if p.color:
                    fill.solid()
                    fill.fore_color.rgb = self._hex_to_rgb(p.color)
                    modified += 1

                elif p.image_path:
                    if not Path(p.image_path).exists():
                        raise FileNotFoundError(f"Background image not found: {p.image_path}")
                    with open(p.image_path, "rb") as f:
                        img_bytes = f.read()
                    # Determine content type
                    ext = Path(p.image_path).suffix.lower()
                    ct_map = {
                        ".png": "image/png",
                        ".jpg": "image/jpeg",
                        ".jpeg": "image/jpeg",
                        ".gif": "image/gif",
                        ".bmp": "image/bmp",
                    }
                    content_type = ct_map.get(ext, "image/png")
                    fill.background()
                    # Use blipFill via XML — python-pptx doesn't have a high-level
                    # API for slide background images; we inject XML directly.
                    from lxml import etree

                    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
                    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"

                    # Add image part to slide
                    from pptx.opc.part import Part
                    from pptx.opc.packuri import PackURI

                    img_part_uri = PackURI(f"/ppt/media/bg_{id(slide)}{ext}")
                    img_part = slide.part.package.part_related_by  # just access the package
                    # Simpler: use slide.part.relate_to with a new blob part
                    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
                    import pptx.opc.part as opc_part

                    image_part = opc_part.Part(
                        img_part_uri, content_type, img_bytes
                    )
                    r_id = slide.part.relate_to(image_part, RT.IMAGE)

                    # Build blipFill XML
                    bg_pr = background._element
                    # Remove existing fill children
                    for child in list(bg_pr):
                        bg_pr.remove(child)

                    bg_fill_xml = (
                        f'<a:blipFill xmlns:a="{ns_a}" xmlns:r="{ns_r}">'
                        f'<a:blip r:embed="{r_id}"/>'
                        f'<a:stretch><a:fillRect/></a:stretch>'
                        f'</a:blipFill>'
                    )
                    bg_fill_elem = etree.fromstring(bg_fill_xml)
                    bg_pr.append(bg_fill_elem)
                    modified += 1

                elif p.gradient:
                    # Gradient fill via XML
                    from lxml import etree

                    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    fill.background()

                    bg_pr = background._element
                    # Remove existing fill children
                    for child in list(bg_pr):
                        bg_pr.remove(child)

                    stops = p.gradient.get("stops", [])
                    angle = p.gradient.get("angle", 0)
                    # OOXML angle is in 60000ths of a degree, clockwise from north
                    ooxml_angle = int(angle * 60000)

                    gs_elems = ""
                    for stop in stops:
                        pos = int(float(stop.get("position", 0)) * 100000)
                        color_hex = stop.get("color", "FFFFFF").lstrip("#")
                        gs_elems += (
                            f'<a:gs pos="{pos}">'
                            f'<a:srgbClr val="{color_hex}"/>'
                            f"</a:gs>"
                        )

                    grad_xml = (
                        f'<a:gradFill xmlns:a="{ns_a}">'
                        f"<a:gsLst>{gs_elems}</a:gsLst>"
                        f'<a:lin ang="{ooxml_angle}" scaled="0"/>'
                        f"</a:gradFill>"
                    )
                    grad_elem = etree.fromstring(grad_xml)
                    bg_pr.append(grad_elem)
                    modified += 1

            self._save_prs(p.path, prs)
            return {"modified_slides": modified}

    @requires_permission(Permission.FILESYSTEM_WRITE, reason="Modifies PowerPoint presentation")
    async def _action_apply_theme(self, params: dict[str, Any]) -> dict[str, Any]:
        p = ApplyThemeParams.model_validate(params)
        return await asyncio.to_thread(self._sync_apply_theme, p)

    def _sync_apply_theme(self, p: ApplyThemeParams) -> dict[str, Any]:
        from pptx import Presentation

        if not Path(p.theme_path).exists():
            raise FileNotFoundError(f"Theme file not found: {p.theme_path}")

        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)
            self._copy_theme_xml(p.theme_path, prs)

            self._save_prs(p.path, prs)
            return {"applied_from": p.theme_path}

    def _copy_theme_xml(self, source_path: str, target_prs: Any) -> None:
        """Copy theme XML from source .pptx into target presentation."""
        from pptx import Presentation

        src_prs = Presentation(source_path)
        src_theme = src_prs.slide_master.theme_color_map
        # Access theme XML element from slide master
        try:
            src_theme_elem = src_prs.slide_master.element.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}theme"
            )
            if src_theme_elem is None:
                # Look in slide master part's relationships for the theme part
                for rel in src_prs.slide_master.part.rels.values():
                    if "theme" in rel.reltype.lower():
                        theme_part = rel.target_part
                        import copy as copy_mod

                        src_theme_xml = copy_mod.deepcopy(theme_part._element)
                        # Find target theme part and replace
                        for trel in target_prs.slide_master.part.rels.values():
                            if "theme" in trel.reltype.lower():
                                tgt_part = trel.target_part
                                tgt_part._element.getparent().replace(
                                    tgt_part._element, src_theme_xml
                                )
                                break
                        break
        except Exception:
            # Theme copying is best-effort
            pass

    @requires_permission(Permission.FILESYSTEM_WRITE, reason="Modifies PowerPoint presentation")
    async def _action_add_transition(self, params: dict[str, Any]) -> dict[str, Any]:
        p = AddTransitionParams.model_validate(params)
        return await asyncio.to_thread(self._sync_add_transition, p)

    def _sync_add_transition(self, p: AddTransitionParams) -> dict[str, Any]:
        from lxml import etree

        with self._get_path_lock(p.path):
            prs = self._get_prs(p.path)

            if p.slide_index is not None:
                if p.slide_index >= len(prs.slides):
                    raise IndexError(f"Slide index {p.slide_index} out of range.")
                target_slides = [prs.slides[p.slide_index]]
            else:
                target_slides = list(prs.slides)

            ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
            # Duration in milliseconds for OOXML
            dur_ms = int(p.duration * 1000)

            for slide in target_slides:
                sp_tree = slide._element
                # Remove existing transition if any
                existing = sp_tree.find(f"{{{ns_p}}}transition")
                if existing is not None:
                    sp_tree.remove(existing)

                if p.transition == "none":
                    continue

                trans_elem = etree.SubElement(sp_tree, f"{{{ns_p}}}transition")
                trans_elem.set("dur", str(dur_ms))
                trans_elem.set("advClick", "1" if p.advance_on_click else "0")

                if p.advance_after is not None:
                    trans_elem.set("advTm", str(int(p.advance_after * 1000)))

                # Add the specific transition child element
                tag_name = _TRANSITION_MAP.get(p.transition, "fade")
                if tag_name:
                    child = etree.SubElement(trans_elem, f"{{{ns_p}}}{tag_name}")
                    if p.transition == "push":
                        child.set("dir", "l")
                    elif p.transition == "wipe":
                        child.set("dir", "l")
                    elif p.transition == "split":
                        child.set("orient", "horz")
                        child.set("dir", "out")

            self._save_prs(p.path, prs)
            return {
                "transition": p.transition,
                "modified_slides": len(target_slides),
            }

    # ------------------------------------------------------------------
    # Export
    # ------------------------------------------------------------------

    async def _action_export_to_pdf(self, params: dict[str, Any]) -> dict[str, Any]:
        p = ExportToPdfParams.model_validate(params)
        return await asyncio.to_thread(self._sync_export_to_pdf, p)

    def _sync_export_to_pdf(self, p: ExportToPdfParams) -> dict[str, Any]:
        with self._get_path_lock(p.path):
            if not Path(p.path).exists():
                raise FileNotFoundError(f"Presentation not found: {p.path}")

            # Ensure the presentation is saved to disk first
            if p.path in self._prs_cache:
                self._prs_cache[p.path].save(p.path)

            output_dir = str(Path(p.output_path).parent)
            Path(output_dir).mkdir(parents=True, exist_ok=True)

            lo_bin = shutil.which("libreoffice") or shutil.which("soffice")
            if lo_bin is None:
                raise RuntimeError(
                    "LibreOffice not found. Install it with: apt install libreoffice"
                )

            result = subprocess.run(
                [
                    lo_bin,
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", output_dir,
                    p.path,
                ],
                capture_output=True,
                text=True,
                timeout=120,
            )
            if result.returncode != 0:
                raise RuntimeError(f"LibreOffice PDF export failed: {result.stderr}")

            # LibreOffice names the output after the input filename
            expected = Path(output_dir) / (Path(p.path).stem + ".pdf")
            if expected.exists() and str(expected) != p.output_path:
                expected.rename(p.output_path)

            return {"pdf_path": p.output_path}

    async def _action_export_slide_as_image(self, params: dict[str, Any]) -> dict[str, Any]:
        p = ExportSlideAsImageParams.model_validate(params)
        return await asyncio.to_thread(self._sync_export_slide_as_image, p)

    def _sync_export_slide_as_image(self, p: ExportSlideAsImageParams) -> dict[str, Any]:
        with self._get_path_lock(p.path):
            if not Path(p.path).exists():
                raise FileNotFoundError(f"Presentation not found: {p.path}")

            if p.path in self._prs_cache:
                self._prs_cache[p.path].save(p.path)

            lo_bin = shutil.which("libreoffice") or shutil.which("soffice")
            if lo_bin is None:
                raise RuntimeError("LibreOffice not found.")

            with tempfile.TemporaryDirectory() as tmpdir:
                # Export all slides as PNG images to tmpdir
                result = subprocess.run(
                    [
                        lo_bin,
                        "--headless",
                        "--convert-to", "png",
                        "--outdir", tmpdir,
                        p.path,
                    ],
                    capture_output=True,
                    text=True,
                    timeout=120,
                )
                if result.returncode != 0:
                    raise RuntimeError(f"LibreOffice image export failed: {result.stderr}")

                # LibreOffice exports slides as <name>1.png, <name>2.png, ...
                stem = Path(p.path).stem
                # Slide index is 0-based, LibreOffice filenames are 1-based
                lo_idx = p.slide_index + 1
                exported = Path(tmpdir) / f"{stem}{lo_idx}.png"

                if not exported.exists():
                    # Fallback: check if only one file was exported
                    pngs = sorted(Path(tmpdir).glob("*.png"))
                    if not pngs:
                        raise RuntimeError("No PNG files generated by LibreOffice.")
                    if p.slide_index < len(pngs):
                        exported = pngs[p.slide_index]
                    else:
                        raise IndexError(
                            f"Slide index {p.slide_index} not found in export "
                            f"({len(pngs)} slides exported)."
                        )

                # Resize if width differs from default
                if p.width != 1920:
                    try:
                        from PIL import Image

                        img = Image.open(exported)
                        ratio = p.width / img.width
                        new_height = int(img.height * ratio)
                        img = img.resize((p.width, new_height), Image.LANCZOS)
                        img.save(str(exported))
                    except ImportError:
                        pass  # PIL not available — skip resize

                Path(p.output_path).parent.mkdir(parents=True, exist_ok=True)
                shutil.copy2(str(exported), p.output_path)

            return {"image_path": p.output_path, "slide_index": p.slide_index}

    # ------------------------------------------------------------------
    # Manifest
    # ------------------------------------------------------------------

    def get_manifest(self) -> ModuleManifest:
        return ModuleManifest(
            module_id="powerpoint",
            version="1.0.0",
            description="Full-featured PowerPoint presentation automation using python-pptx.",
            platforms=["all"],
            tags=["powerpoint", "pptx", "presentation", "slides", "office"],
            dependencies=["python-pptx>=0.6"],
            declared_permissions=["filesystem_read", "filesystem_write"],
            actions=[
                # -- Lifecycle --
                ActionSpec(
                    name="create_presentation",
                    description="Create a new blank PowerPoint presentation.",
                    params=[
                        ParamSpec("output_path", "string", "Path where the .pptx will be saved."),
                        ParamSpec("slide_width", "number", "Slide width in cm.", required=False),
                        ParamSpec("slide_height", "number", "Slide height in cm.", required=False),
                        ParamSpec("theme_path", "string", "Source .pptx to copy theme from.", required=False),
                    ],
                    returns="object",
                    returns_description='{"path": str, "slide_width_cm": float, "slide_height_cm": float}',
                ),
                ActionSpec(
                    name="open_presentation",
                    description="Open and cache an existing .pptx file.",
                    params=[
                        ParamSpec("path", "string", "Path to an existing .pptx file."),
                    ],
                    returns="object",
                    returns_description='{"path": str, "slide_count": int, "layout_names": list[str]}',
                ),
                ActionSpec(
                    name="save_presentation",
                    description="Save the presentation to disk.",
                    params=[
                        ParamSpec("path", "string", "Path of the cached presentation."),
                        ParamSpec("output_path", "string", "Save-as path. Overwrites original if omitted.", required=False),
                    ],
                    returns="object",
                    returns_description='{"saved_to": str}',
                ),
                ActionSpec(
                    name="get_presentation_info",
                    description="Return metadata about a presentation.",
                    params=[ParamSpec("path", "string", "Path to the .pptx file.")],
                    returns="object",
                    returns_description='{"slide_count": int, "slide_width_cm": float, "slide_height_cm": float, "slide_layouts": list[str]}',
                ),
                # -- Slide management --
                ActionSpec(
                    name="add_slide",
                    description="Add a new slide to the presentation.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("layout_index", "integer", "Slide layout index.", required=False, default=1),
                        ParamSpec("title", "string", "Optional slide title.", required=False),
                        ParamSpec("position", "integer", "Insert position (0-indexed). Appended if omitted.", required=False),
                    ],
                    returns="object",
                    returns_description='{"slide_index": int, "slide_count": int}',
                ),
                ActionSpec(
                    name="delete_slide",
                    description="Delete a slide from the presentation.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("slide_index", "integer", "0-indexed slide to delete."),
                    ],
                    returns="object",
                ),
                ActionSpec(
                    name="duplicate_slide",
                    description="Duplicate a slide and insert the copy at a specified position.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("slide_index", "integer", "0-indexed slide to duplicate."),
                        ParamSpec("insert_after", "integer", "Insert copy after this index.", required=False),
                    ],
                    returns="object",
                ),
                ActionSpec(
                    name="reorder_slide",
                    description="Move a slide from one position to another.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("from_index", "integer", "Source slide index."),
                        ParamSpec("to_index", "integer", "Destination slide index."),
                    ],
                    returns="object",
                ),
                ActionSpec(
                    name="list_slides",
                    description="List all slides with title, shape count, and notes preview.",
                    params=[ParamSpec("path", "string", "Path to the .pptx file.")],
                    returns="object",
                    returns_description='{"slides": [{"index": int, "title": str, "shape_count": int, "notes_preview": str}]}',
                ),
                ActionSpec(
                    name="read_slide",
                    description="Read all content from a slide.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("slide_index", "integer", "0-indexed slide to read."),
                        ParamSpec("include_notes", "boolean", "Include speaker notes.", required=False, default=True),
                        ParamSpec("include_shapes", "boolean", "Include shape details.", required=False, default=True),
                    ],
                    returns="object",
                ),
                ActionSpec(
                    name="set_slide_layout",
                    description="Change the layout of a slide.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("slide_index", "integer", "0-indexed slide."),
                        ParamSpec("layout_index", "integer", "Target layout index."),
                    ],
                    returns="object",
                ),
                # -- Text content --
                ActionSpec(
                    name="set_slide_title",
                    description="Set the title placeholder text of a slide.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("slide_index", "integer", "0-indexed slide."),
                        ParamSpec("title", "string", "Title text."),
                        ParamSpec("bold", "boolean", "Bold title text.", required=False, default=False),
                        ParamSpec("font_size", "integer", "Font size in points.", required=False),
                        ParamSpec("font_color", "string", "Hex colour, e.g. FF0000.", required=False),
                    ],
                    returns="object",
                ),
                ActionSpec(
                    name="add_text_box",
                    description="Add a text box to a slide with full formatting options.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("slide_index", "integer", "0-indexed slide."),
                        ParamSpec("text", "string", "Text content."),
                        ParamSpec("left", "number", "Left position in cm."),
                        ParamSpec("top", "number", "Top position in cm."),
                        ParamSpec("width", "number", "Width in cm."),
                        ParamSpec("height", "number", "Height in cm."),
                        ParamSpec("bold", "boolean", "Bold text.", required=False, default=False),
                        ParamSpec("italic", "boolean", "Italic text.", required=False, default=False),
                        ParamSpec("font_size", "integer", "Font size in pt.", required=False),
                        ParamSpec("font_color", "string", "Hex colour.", required=False),
                        ParamSpec("alignment", "string", "Text alignment: left|center|right|justify.", required=False, default="left"),
                    ],
                    returns="object",
                ),
                ActionSpec(
                    name="add_slide_notes",
                    description="Set or append speaker notes to a slide.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("slide_index", "integer", "0-indexed slide."),
                        ParamSpec("notes", "string", "Notes text."),
                        ParamSpec("append", "boolean", "Append to existing notes.", required=False, default=False),
                    ],
                    returns="object",
                ),
                # -- Shapes --
                ActionSpec(
                    name="add_shape",
                    description="Add an auto shape to a slide.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("slide_index", "integer", "0-indexed slide."),
                        ParamSpec("shape_type", "string", "Shape type: rectangle|ellipse|triangle|...", required=False, default="rectangle"),
                        ParamSpec("left", "number", "Left position in cm."),
                        ParamSpec("top", "number", "Top position in cm."),
                        ParamSpec("width", "number", "Width in cm."),
                        ParamSpec("height", "number", "Height in cm."),
                        ParamSpec("fill_color", "string", "Fill hex colour.", required=False),
                        ParamSpec("line_color", "string", "Border hex colour.", required=False),
                        ParamSpec("text", "string", "Text inside shape.", required=False),
                    ],
                    returns="object",
                ),
                ActionSpec(
                    name="format_shape",
                    description="Modify fill, border, rotation, and shadow of an existing shape.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("slide_index", "integer", "0-indexed slide."),
                        ParamSpec("shape_index", "integer", "0-indexed shape."),
                        ParamSpec("fill_color", "string", "Fill hex colour.", required=False),
                        ParamSpec("line_color", "string", "Border hex colour.", required=False),
                        ParamSpec("rotation", "number", "Rotation degrees.", required=False),
                        ParamSpec("shadow", "boolean", "Apply drop shadow.", required=False),
                    ],
                    returns="object",
                ),
                ActionSpec(
                    name="add_image",
                    description="Insert an image onto a slide.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("slide_index", "integer", "0-indexed slide."),
                        ParamSpec("image_path", "string", "Path to the image file."),
                        ParamSpec("left", "number", "Left position in cm."),
                        ParamSpec("top", "number", "Top position in cm."),
                        ParamSpec("width", "number", "Width in cm. Auto-scaled if omitted.", required=False),
                        ParamSpec("height", "number", "Height in cm. Auto-scaled if omitted.", required=False),
                    ],
                    returns="object",
                ),
                # -- Charts --
                ActionSpec(
                    name="add_chart",
                    description="Add a chart to a slide.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("slide_index", "integer", "0-indexed slide."),
                        ParamSpec("chart_type", "string", "Chart type: bar|col|line|pie|doughnut|scatter|area|bubble|radar.", required=False, default="col"),
                        ParamSpec("data", "object", 'Chart data: {"categories": [...], "series": [{"name": str, "values": [...]}]}'),
                        ParamSpec("left", "number", "Left position in cm."),
                        ParamSpec("top", "number", "Top position in cm."),
                        ParamSpec("width", "number", "Width in cm.", required=False, default=14),
                        ParamSpec("height", "number", "Height in cm.", required=False, default=10),
                        ParamSpec("title", "string", "Chart title.", required=False),
                        ParamSpec("has_legend", "boolean", "Show legend.", required=False, default=True),
                        ParamSpec("has_data_labels", "boolean", "Show data labels.", required=False, default=False),
                    ],
                    returns="object",
                ),
                # -- Tables --
                ActionSpec(
                    name="add_table",
                    description="Add a table to a slide.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("slide_index", "integer", "0-indexed slide."),
                        ParamSpec("rows", "integer", "Number of rows."),
                        ParamSpec("cols", "integer", "Number of columns."),
                        ParamSpec("data", "array", "Row-major list of cell values.", required=False),
                        ParamSpec("left", "number", "Left position in cm."),
                        ParamSpec("top", "number", "Top position in cm."),
                        ParamSpec("width", "number", "Table width in cm.", required=False, default=20),
                        ParamSpec("height", "number", "Table height in cm.", required=False, default=10),
                        ParamSpec("has_header", "boolean", "Style first row as header.", required=False, default=True),
                    ],
                    returns="object",
                ),
                ActionSpec(
                    name="format_table_cell",
                    description="Format a specific cell in a table shape.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("slide_index", "integer", "0-indexed slide."),
                        ParamSpec("shape_index", "integer", "0-indexed table shape."),
                        ParamSpec("row", "integer", "Row index (0-based)."),
                        ParamSpec("col", "integer", "Column index (0-based)."),
                        ParamSpec("text", "string", "Cell text.", required=False),
                        ParamSpec("bg_color", "string", "Background hex colour.", required=False),
                        ParamSpec("font_color", "string", "Font hex colour.", required=False),
                        ParamSpec("bold", "boolean", "Bold text.", required=False),
                    ],
                    returns="object",
                ),
                # -- Background & theme --
                ActionSpec(
                    name="set_slide_background",
                    description="Set the background of one or all slides (solid color, image, or gradient).",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("slide_index", "integer", "0-indexed slide. All slides if omitted.", required=False),
                        ParamSpec("color", "string", "Solid background hex colour.", required=False),
                        ParamSpec("image_path", "string", "Path to background image.", required=False),
                        ParamSpec("gradient", "object", "Gradient spec: {type, stops, angle}.", required=False),
                    ],
                    returns="object",
                ),
                ActionSpec(
                    name="apply_theme",
                    description="Copy and apply the theme from another .pptx file.",
                    params=[
                        ParamSpec("path", "string", "Path to the target .pptx."),
                        ParamSpec("theme_path", "string", "Path to the source .pptx to copy theme from."),
                    ],
                    returns="object",
                ),
                ActionSpec(
                    name="add_transition",
                    description="Add a slide transition animation.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("slide_index", "integer", "0-indexed slide. All slides if omitted.", required=False),
                        ParamSpec("transition", "string", "Transition type: none|fade|push|wipe|split|reveal|random.", required=False, default="fade"),
                        ParamSpec("duration", "number", "Transition duration in seconds.", required=False, default=1.0),
                        ParamSpec("advance_on_click", "boolean", "Advance on click.", required=False, default=True),
                        ParamSpec("advance_after", "number", "Auto-advance after N seconds.", required=False),
                    ],
                    returns="object",
                ),
                # -- Export --
                ActionSpec(
                    name="export_to_pdf",
                    description="Export the presentation to a PDF file using LibreOffice.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("output_path", "string", "Destination PDF path."),
                    ],
                    returns="object",
                    returns_description='{"pdf_path": str}',
                ),
                ActionSpec(
                    name="export_slide_as_image",
                    description="Export a single slide as a PNG image using LibreOffice.",
                    params=[
                        ParamSpec("path", "string", "Path to the .pptx file."),
                        ParamSpec("slide_index", "integer", "0-indexed slide to export."),
                        ParamSpec("output_path", "string", "Destination image path."),
                        ParamSpec("width", "integer", "Output width in pixels.", required=False, default=1920),
                    ],
                    returns="object",
                    returns_description='{"image_path": str, "slide_index": int}',
                ),
            ],
        )
